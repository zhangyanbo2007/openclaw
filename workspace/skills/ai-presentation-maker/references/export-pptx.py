#!/usr/bin/env python3
"""
export-pptx.py — Convert a Presentation Maker markdown deck to PPTX.
Requires: python-pptx (pip install python-pptx)

Usage:
  python3 export-pptx.py <deck.md> <output.pptx> [metadata.json]
"""

import sys
import re
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def parse_markdown_slides(md_path: str) -> list:
    """Parse markdown file into slide objects."""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []
    current_slide = None

    for line in content.split('\n'):
        # Title slide (# heading)
        if line.startswith('# ') and not line.startswith('## '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'title',
                'title': line[2:].strip(),
                'content': [],
                'speaker_notes': []
            }
        # Regular slide (## heading)
        elif line.startswith('## '):
            if current_slide:
                slides.append(current_slide)
            # Extract slide title (remove "Slide N: " prefix if present)
            title = line[3:].strip()
            title = re.sub(r'^Slide \d+:\s*', '', title)
            current_slide = {
                'type': 'content',
                'title': title,
                'content': [],
                'speaker_notes': []
            }
        # Speaker notes marker
        elif current_slide and ('**Speaker Notes' in line or '**Say:**' in line or '**What to say' in line):
            current_slide['_in_notes'] = True
        elif current_slide and line.startswith('---'):
            if current_slide.get('_in_notes'):
                current_slide.pop('_in_notes', None)
        elif current_slide:
            if current_slide.get('_in_notes'):
                # Clean markdown formatting for notes
                clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', line.strip())
                clean = re.sub(r'^>\s*', '', clean)
                if clean:
                    current_slide['speaker_notes'].append(clean)
            else:
                clean = line.strip()
                if clean and not clean.startswith('>'):
                    current_slide['content'].append(clean)

    if current_slide:
        slides.append(current_slide)

    # Clean up internal flags
    for s in slides:
        s.pop('_in_notes', None)

    return slides


def create_pptx(slides: list, output_path: str, metadata: dict = None):
    """Generate a PPTX file from parsed slides."""
    prs = Presentation()

    # Set slide dimensions (widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        if slide_data['type'] == 'title':
            # Title slide
            layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(layout)

            title = slide.shapes.title
            if title:
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.size = Pt(36)

            # Add subtitle from metadata
            if metadata and slide.placeholders:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:  # Subtitle
                        speaker = metadata.get('speaker', {})
                        subtitle_parts = []
                        if speaker.get('name'):
                            subtitle_parts.append(speaker['name'])
                        if speaker.get('title'):
                            subtitle_parts.append(speaker['title'])
                        ph.text = ' — '.join(subtitle_parts)
                        break

        else:
            # Content slide
            layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(layout)

            # Set title
            title = slide.shapes.title
            if title:
                title.text = slide_data['title']
                title.text_frame.paragraphs[0].font.size = Pt(28)

            # Set content
            content_text = '\n'.join(slide_data['content'])
            # Clean markdown formatting
            content_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', content_text)
            content_text = re.sub(r'\*([^*]+)\*', r'\1', content_text)
            content_text = re.sub(r'^- ', '• ', content_text, flags=re.MULTILINE)
            content_text = re.sub(r'^\d+\. ', '• ', content_text, flags=re.MULTILINE)

            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # Content area
                    ph.text = content_text
                    for para in ph.text_frame.paragraphs:
                        para.font.size = Pt(18)
                    break

        # Add speaker notes
        if slide_data['speaker_notes']:
            notes_text = '\n'.join(slide_data['speaker_notes'])
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    prs.save(output_path)
    print(f"✅ PPTX saved: {output_path} ({len(slides)} slides)")


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 export-pptx.py <deck.md> <output.pptx> [metadata.json]")
        sys.exit(1)

    md_path = sys.argv[1]
    output_path = sys.argv[2]
    meta_path = sys.argv[3] if len(sys.argv) > 3 else None

    metadata = None
    if meta_path and Path(meta_path).exists():
        with open(meta_path, 'r', encoding='utf-8') as f:
            metadata = json.load(f)

    slides = parse_markdown_slides(md_path)
    if not slides:
        print("ERROR: No slides found in markdown file", file=sys.stderr)
        sys.exit(1)

    create_pptx(slides, output_path, metadata)


if __name__ == '__main__':
    main()
